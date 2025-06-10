import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_ppt_advanced(file_path):
    """
    Extracts text from a PowerPoint, preserving structure like titles,
    table cells, and bold formatting.
    """
    if not os.path.exists(file_path):
        print(f"Error: The file '{file_path}' was not found.")
        return []

    print(f"Opening presentation for deep extraction: {file_path}")
    prs = Presentation(file_path)
    text_map = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            text_map.append({
                                'type': 'table_cell',
                                'slide_index': slide_index,
                                'shape_index': shape_index,
                                'location': {'row': r_idx, 'col': c_idx},
                                'original_text': cell.text.strip(),
                                'is_bold': False
                            })
            elif shape.has_text_frame:
                is_title = (shape.is_placeholder and
                            hasattr(shape.placeholder_format, 'type') and
                            'TITLE' in str(shape.placeholder_format.type))
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            text_map.append({
                                'type': 'text_run',
                                'is_title': is_title,
                                'slide_index': slide_index,
                                'shape_index': shape_index,
                                'location': {'paragraph': p_idx, 'run': r_idx},
                                'original_text': run.text.strip(),
                                'is_bold': run.font.bold or False
                            })
    return text_map

# This block allows us to test this file directly
if __name__ == '__main__':
    print("Testing extractor.py directly...")
    test_data = extract_text_from_ppt_advanced('sample_deck.pptx')
    if test_data:
        print(f"Successfully extracted {len(test_data)} elements.")
        print("Sample item:", test_data[0])
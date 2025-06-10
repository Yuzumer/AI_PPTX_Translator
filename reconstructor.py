from pptx import Presentation


def reconstruct_presentation(text_map, original_ppt_path, output_ppt_path):
    """
    Reconstructs a presentation with translated text, preserving formatting,
    with corrected, robust handling for table cell text.
    """
    print("\n--- Starting Phase 4: Reconstructing the Presentation ---")

    try:
        prs = Presentation(original_ppt_path)
    except Exception as e:
        print(f"Error opening original presentation file: {e}")
        return

    print(f"Updating {len(text_map)} text elements...")

    for item in text_map:
        try:
            slide = prs.slides[item['slide_index']]
            shape = slide.shapes[item['shape_index']]
            translated_text = item.get('translated_text', item['original_text'])

            if item['type'] == 'table_cell':
                row, col = item['location']['row'], item['location']['col']
                cell = shape.table.cell(row, col)

                # --- BUG FIX STARTS HERE ---

                text_frame = cell.text_frame

                # First, ensure the cell has at least one paragraph to work with.
                if not text_frame.paragraphs:
                    # If the cell is completely empty, add a paragraph.
                    p = text_frame.add_paragraph()
                else:
                    # THE FIX: We must explicitly target the FIRST paragraph (index 0).
                    # The old code incorrectly tried to operate on the entire list of paragraphs.
                    p = text_frame.paragraphs[0]

                # Now that 'p' is a single paragraph object, we can safely modify it.
                # This simple assignment is the most reliable way to set text and preserve
                # the dominant style of the paragraph.
                p.text = translated_text

                # Clean-up: If there were multiple paragraphs in the cell before,
                # remove them to ensure only the new translated text remains.
                # This loop safely removes paragraphs from the end backwards.
                for i in range(len(text_frame.paragraphs) - 1, 0, -1):
                    p_to_remove = text_frame.paragraphs[i]
                    text_frame._txBody.remove(p_to_remove._p)

                # --- BUG FIX ENDS HERE ---

            elif item['type'] == 'text_run':
                p_idx, r_idx = item['location']['paragraph'], item['location']['run']
                run = shape.text_frame.paragraphs[p_idx].runs[r_idx]
                run.text = translated_text

                if item.get('is_bold'):
                    run.font.bold = True
                else:
                    run.font.bold = False

        except (IndexError, KeyError) as e:
            print(
                f"  - Warning: Could not find or process element at {item.get('location', 'N/A')}. Error: {e}. Skipping.")
        except Exception as e:
            # This is where our 'tuple' error was being caught.
            print(f"  - Warning: An unexpected error occurred while updating an item: {e}. Skipping.")

    try:
        print(f"\nSaving translated presentation to: {output_ppt_path}")
        prs.save(output_ppt_path)
        print("--- Reconstruction Complete ---")
    except Exception as e:
        print(f"Error saving the final presentation: {e}")